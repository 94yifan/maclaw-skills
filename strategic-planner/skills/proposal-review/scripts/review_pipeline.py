#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
proposal-review pipeline —— 方案审查流水线（固化版）

把 proposal-review SKILL 的审查 SOP 固化成一个可重复执行的 pipeline：
  1) 前置 gate：先要 brief —— 没有 brief 直接拦停，输出生意盘子三问
  2) 内容抽取：PPTX / PDF / 纯文本 → 分页 / 分板块文本
  3) 检查清单：按 Step 0-10 生成完整审查工作表（含产品层新检查）
  4) 报告模板：输出统一结构的审查报告骨架

审查判断仍由 Pro 模型完成；本脚本负责「抽取 + 前置拦截 + 清单 + 模板」，
保证每一份方案走同一条审查流水线，不漏步骤、不跳前置。

用法:
  python review_pipeline.py --file proposal.pdf --brief brief.txt
  python review_pipeline.py --file proposal.pptx --brief "中秋礼盒单品战役；目标提升特别特搜索指数"
  python review_pipeline.py --text extracted.txt --brief brief.txt
  python review_pipeline.py --file proposal.pdf          # 缺 brief → 拦停输出三问
  python review_pipeline.py --file proposal.pdf --brief b.txt --out worksheet.md

依赖（按需）:
  - .pptx 抽取需要 python-pptx：pip install python-pptx
  - .pdf 文本抽取需要 pdftotext（poppler-utils）；图片型 fallback 需要 pdftoppm
"""

import argparse
import os
import subprocess
import sys
import tempfile


THREE_QUESTIONS = """\
审查前先确认「生意盘子三问」（brief 应能回答）：
  1. 客户委托的盘子是什么：系列战略 / 单品战役 / 节点 / 执行
  2. 这个产品的生意本质是什么：礼盒=礼赠、节令=节点、速食=即得
  3. 曼拾这次该交付的核心价值是什么：投流统筹 / KFS 统筹 / 达人

三问答不出 = 审查失焦。没有 brief 就先向逸凡要，别凭方案倒推猜产品。
"""


REPORT_TEMPLATE = """\
## 审查前置（Step 0）
Brief：[有 / 无，无则先列三问待确认]
生意盘子三问：
  1. 委托盘子：
  2. 产品生意本质：
  3. 曼拾交付价值：
产品五查（立不立得住）：卖点独特 / KM独有资产 / 零理解成本 / 层级 / 词根延展
产品穿透链：卖点→人群(自用+礼赠)→内容→脚本→投放，对象是否一致
礼赠/节令：礼赠人群是否含 / 时间轴是否对齐节日

## 审查概要
[方案名] | [格式] | [方案类型] | [总规模] | [日期]
适用标准：策划方案框架 + 前策分析培训文档
核心问题：[1-3个]

## 分段控制：强因果区 / 弱因果区
强因果区：[范围]
弱因果区：[范围]

## 模块零：方案类型匹配
## 模块一：结构逻辑（含平台策略页完整性）
## 模块二：策略推导 + 段落内部逻辑（含一句话策略+品类核心矛盾）
## 模块三：数据来源合规
## 模块四：KFS预算分配 + 品牌阶段战略匹配（含预算配置合理性）
## 模块五：因果链（含产品穿透链）
## 模块六：数据质量 + 颗粒度（含礼赠/节令规则）
## 模块七：矩阵与内容（含KOL匹配产品特性+脚本专业化）
## 模块八：修改追踪 + 连锁影响

## 修改优先级汇总
P0：
P1：
P2：

--------------------------------------------------------------------
## 增量建议（非必要，供参考）
（三段式：数据线索 + 机会推导 + 内容建议）
--------------------------------------------------------------------

审查完成。
"""


def fail(msg, code=2):
    print(msg, file=sys.stderr)
    sys.exit(code)


def read_text_arg(arg):
    """brief/text 参数既可能是文件路径，也可能是内联文本。"""
    if os.path.isfile(arg):
        with open(arg, encoding="utf-8") as f:
            return f.read().strip()
    return arg.strip()


def extract_pptx(path):
    try:
        from pptx import Presentation
    except ImportError:
        fail("缺少 python-pptx，请先安装：pip install python-pptx")
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        out.append("\n===== 第%d页 =====\n" % i + "\n".join(texts))
    return "\n".join(out)


def extract_pdf(path):
    # 文本型：pdftotext；图片型：pdftoppm 转图片，交给 image 工具逐页读
    if subprocess.run(["which", "pdftotext"], capture_output=True).returncode == 0:
        r = subprocess.run(
            ["pdftotext", "-layout", path, "-"], capture_output=True, text=True
        )
        txt = r.stdout.strip()
        if len(txt) > 200:
            return txt
    if subprocess.run(["which", "pdftoppm"], capture_output=True).returncode == 0:
        tmp = tempfile.mkdtemp(prefix="pr_pages_")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "120", path, os.path.join(tmp, "p")],
            capture_output=True,
        )
        pages = sorted(os.listdir(tmp))
        lines = [
            "[图片型PDF] 第%d页 -> %s" % (i + 1, os.path.join(tmp, p))
            for i, p in enumerate(pages)
        ]
        return (
            "\n".join(lines)
            + "\n\n(共 %d 页，图片型 PDF，需用 image 工具逐页读)" % len(pages)
        )
    fail("PDF 抽取失败：需要 pdftotext 或 pdftoppm（poppler-utils）")


def detect_and_extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext in (".pptx", ".ppt"):
        return "pptx", extract_pptx(path)
    if ext == ".pdf":
        return "pdf", extract_pdf(path)
    if ext in (".txt", ".md"):
        with open(path, encoding="utf-8") as f:
            return "text", f.read()
    fail("不支持的格式: %s（支持 .pptx / .pdf / .txt / .md）" % ext)


def build_worksheet(fmt, brief, content, src):
    brief_preview = brief if len(brief) <= 300 else brief[:300] + "..."
    return (
        "# 方案审查 Worksheet（pipeline 生成）\n\n"
        "## 输入\n"
        "- 方案：%s\n"
        "- 格式：%s\n"
        "- Brief：%s\n\n"
        "%s\n\n"
        "---\n\n"
        "## 方案正文（逐页/逐板块对照审查）\n\n%s\n"
    ) % (src, fmt, brief_preview, REPORT_TEMPLATE, content[:60000])


def main():
    ap = argparse.ArgumentParser(
        description="proposal-review pipeline：前置 gate + 内容抽取 + 检查清单 + 报告模板"
    )
    ap.add_argument("--file", help="方案文件路径（.pptx / .pdf / .txt / .md）")
    ap.add_argument("--text", help="已抽取的纯文本文件路径（飞书 / 手动场景）")
    ap.add_argument("--brief", default="", help="客户委托 brief（文件路径或内联文本）")
    ap.add_argument("--out", help="输出 worksheet 到文件（默认 stdout）")
    args = ap.parse_args()

    # ---------- 前置 gate：先要 brief ----------
    if not args.brief:
        print("=" * 62)
        print("STEP 0 前置拦截：缺 brief")
        print("=" * 62)
        print(THREE_QUESTIONS)
        print("请补充 brief 后重跑：--brief <内容或文件路径>")
        print("brief 是审查基准，没有它产品对不对、目标对不对都无从判断。")
        sys.exit(2)

    brief = read_text_arg(args.brief)

    # ---------- 内容抽取 ----------
    if args.text:
        content = read_text_arg(args.text)
        fmt = "text"
        src = args.text
    elif args.file:
        fmt, content = detect_and_extract(args.file)
        src = args.file
    else:
        fail("需提供 --file 或 --text 之一")

    # ---------- 生成 worksheet ----------
    worksheet = build_worksheet(fmt, brief, content, src)

    if args.out:
        with open(args.out, "w", encoding="utf-8") as f:
            f.write(worksheet)
        print("worksheet 已写入 %s" % args.out)
    else:
        print(worksheet)


if __name__ == "__main__":
    main()
